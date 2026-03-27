#!/usr/bin/env python3
"""
Generate a macro-enabled Jeopardy PowerPoint from a spreadsheet.

Usage:
    python3 generate_jeopardy_pptm.py --input sample_data/sample_jeopardy_board.xlsx --output build/jeopardy_game.pptm
"""

from __future__ import annotations

import argparse
import math
import re
import shutil
import subprocess
import sys
import textwrap
import zipfile
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

COLOR_BG = RGBColor(7, 16, 64)
COLOR_BOARD = RGBColor(13, 53, 176)
COLOR_BOARD_USED = RGBColor(33, 42, 76)
COLOR_HEADER = RGBColor(24, 74, 208)
COLOR_PANEL = RGBColor(16, 34, 109)
COLOR_BUTTON = RGBColor(255, 198, 58)
COLOR_TEXT = RGBColor(255, 255, 255)
COLOR_ACCENT = RGBColor(118, 221, 255)


@dataclass(frozen=True)
class ClueRecord:
    clue_id: str
    safe_id: str
    ordinal: int
    category: str
    category_sort: tuple
    value_label: str
    value_sort: tuple
    clue: str
    answer: str


@dataclass(frozen=True)
class SlideMeta:
    clue_index: int
    answer_index: int


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Generate a macro-enabled Jeopardy deck.")
    parser.add_argument("--input", required=True, help="Input .xlsx or .csv file")
    parser.add_argument("--output", required=True, help="Output .pptm path")
    parser.add_argument("--report", default="build/jeopardy_macro_report.txt", help="Validation report path")
    return parser.parse_args()


def normalize_name(name: object) -> str:
    return re.sub(r"[^a-z0-9]+", "", ("" if name is None else str(name)).strip().lower())


def safe_vba_id(text: str) -> str:
    clean = re.sub(r"[^A-Za-z0-9_]+", "_", text)
    clean = re.sub(r"_+", "_", clean).strip("_")
    if not clean:
        clean = "clue"
    if clean[0].isdigit():
        clean = f"id_{clean}"
    return clean


def best_column_match(columns: Iterable[object], aliases: list[str], required: bool = True) -> object | None:
    normalized = {column: normalize_name(column) for column in columns}
    alias_set = {normalize_name(alias) for alias in aliases}

    for column, value in normalized.items():
        if value in alias_set:
            return column

    for column, value in normalized.items():
        if any(alias in value or value in alias for alias in alias_set):
            return column

    import difflib

    matches: list[tuple[float, object]] = []
    for column, value in normalized.items():
        score = max(difflib.SequenceMatcher(None, value, alias).ratio() for alias in alias_set)
        matches.append((score, column))

    if not matches:
        if required:
            raise ValueError("No columns available for matching.")
        return None

    matches.sort(reverse=True, key=lambda item: item[0])
    score, column = matches[0]
    if required and score < 0.50:
        raise ValueError(f"Could not confidently match any of {aliases} in columns {list(columns)}")
    return column if score >= 0.35 else None


def load_frame(input_path: Path) -> tuple[pd.DataFrame, str]:
    if input_path.suffix.lower() == ".csv":
        return pd.read_csv(input_path), input_path.name

    excel = pd.ExcelFile(input_path)
    required = {
        "category": ["category", "categories", "topic"],
        "value": ["value", "points", "amount", "dollarvalue", "score"],
        "clue": ["clue", "question", "prompt"],
        "answer": ["answer", "response", "expectedresponse", "expected_response", "correctresponse"],
    }

    best: tuple[int, str, pd.DataFrame] | None = None
    for sheet_name in excel.sheet_names:
        frame = excel.parse(sheet_name=sheet_name).dropna(how="all")
        if frame.empty:
            continue
        score = 0
        for aliases in required.values():
            try:
                best_column_match(frame.columns, aliases, required=True)
                score += 1
            except ValueError:
                pass
        if best is None or score > best[0]:
            best = (score, sheet_name, frame)

    if best is None or best[0] < 4:
        raise ValueError("Could not find a usable worksheet.")
    return best[2], best[1]


def coerce_text(value: object) -> str:
    if value is None or (isinstance(value, float) and math.isnan(value)):
        return ""
    return str(value).strip()


def parse_numeric_value(value: object) -> float | None:
    text = coerce_text(value)
    if not text:
        return None
    match = re.search(r"-?\d+(?:\.\d+)?", text.replace(",", ""))
    return float(match.group()) if match else None


def load_records(input_path: Path) -> tuple[list[ClueRecord], str]:
    frame, source_name = load_frame(input_path)
    frame = frame.dropna(how="all").copy()

    cols = {
        "category": best_column_match(frame.columns, ["category", "categories", "topic"]),
        "value": best_column_match(frame.columns, ["value", "points", "amount", "score"]),
        "clue": best_column_match(frame.columns, ["clue", "question", "prompt"]),
        "answer": best_column_match(frame.columns, ["answer", "response", "expectedresponse", "expected_response"]),
        "category_order": best_column_match(frame.columns, ["category_order", "categoryorder"], required=False),
        "value_order": best_column_match(frame.columns, ["value_order", "valueorder"], required=False),
        "id": best_column_match(frame.columns, ["id", "clue_id", "question_id"], required=False),
    }

    raw_rows: list[dict[str, object]] = []
    for row_index, row in frame.iterrows():
        category = coerce_text(row[cols["category"]])
        value_label = coerce_text(row[cols["value"]])
        clue = coerce_text(row[cols["clue"]])
        answer = coerce_text(row[cols["answer"]])
        if not (category and value_label and clue and answer):
            continue
        raw_rows.append(
            {
                "row_number": row_index + 2,
                "clue_id": coerce_text(row[cols["id"]]) if cols["id"] is not None else f"row-{row_index + 2}",
                "category": category,
                "category_order": row[cols["category_order"]] if cols["category_order"] is not None else None,
                "value_label": value_label,
                "value_order": row[cols["value_order"]] if cols["value_order"] is not None else None,
                "value_number": parse_numeric_value(row[cols["value"]]),
                "clue": clue,
                "answer": answer,
            }
        )

    if not raw_rows:
        raise ValueError("No complete clue rows were found.")

    duplicates: dict[tuple[str, str], list[int]] = {}
    for row in raw_rows:
        key = (str(row["category"]), str(row["value_label"]))
        duplicates.setdefault(key, []).append(int(row["row_number"]))
    duplicate_hits = {k: v for k, v in duplicates.items() if len(v) > 1}
    if duplicate_hits:
        raise ValueError(f"Duplicate category/value pairs found: {duplicate_hits}")

    category_positions: dict[str, int] = {}
    value_positions: dict[str, int] = {}
    for row in raw_rows:
        category_positions.setdefault(str(row["category"]), len(category_positions))
        value_positions.setdefault(str(row["value_label"]), len(value_positions))

    raw_rows.sort(
        key=lambda row: (
            float(row["value_order"]) if pd.notna(row["value_order"]) else float("inf"),
            float(row["value_number"]) if row["value_number"] is not None else float("inf"),
            value_positions[str(row["value_label"])],
            float(row["category_order"]) if pd.notna(row["category_order"]) else float("inf"),
            category_positions[str(row["category"])],
        )
    )

    records: list[ClueRecord] = []
    for ordinal, row in enumerate(raw_rows, start=1):
        category = str(row["category"])
        value_label = str(row["value_label"])
        records.append(
            ClueRecord(
                clue_id=str(row["clue_id"]),
                safe_id=safe_vba_id(str(row["clue_id"])),
                ordinal=ordinal,
                category=category,
                category_sort=(
                    float(row["category_order"]) if pd.notna(row["category_order"]) else float("inf"),
                    category_positions[category],
                    category.lower(),
                ),
                value_label=value_label,
                value_sort=(
                    float(row["value_order"]) if pd.notna(row["value_order"]) else float("inf"),
                    float(row["value_number"]) if row["value_number"] is not None else float("inf"),
                    value_positions[value_label],
                    value_label.lower(),
                ),
                clue=str(row["clue"]),
                answer=str(row["answer"]),
            )
        )
    return records, source_name


def order_categories(records: list[ClueRecord]) -> list[str]:
    categories = {record.category: record.category_sort for record in records}
    return [name for name, _ in sorted(categories.items(), key=lambda item: item[1])]


def order_values(records: list[ClueRecord]) -> list[str]:
    values = {record.value_label: record.value_sort for record in records}
    return [name for name, _ in sorted(values.items(), key=lambda item: item[1])]


def set_shape_name(shape, name: str) -> None:
    shape._element.nvSpPr.cNvPr.set("name", name)


def set_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG


def fit_font_size(text: str, *, max_chars_per_line: int, max_font: int, min_font: int) -> int:
    line_estimate = max(1, math.ceil(len(text) / max_chars_per_line))
    if line_estimate <= 2:
        return max_font
    return max(min_font, max_font - ((line_estimate - 2) * 2))


def set_shape_text(shape, text: str, *, font_size: int, bold: bool, color: RGBColor,
                   align=PP_ALIGN.CENTER, vertical=MSO_ANCHOR.MIDDLE) -> None:
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = vertical
    frame.margin_left = Pt(6)
    frame.margin_right = Pt(6)
    frame.margin_top = Pt(6)
    frame.margin_bottom = Pt(6)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_rect(slide, left, top, width, height, *, fill_color: RGBColor, line_color: RGBColor, line_width: float = 1.2):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color
    shape.line.width = Pt(line_width)
    return shape


def create_board_slide(prs: Presentation, categories: list[str], values: list[str], records_by_key: dict[tuple[str, str], ClueRecord]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)

    title = slide.shapes.add_textbox(Inches(0.32), Inches(0.12), Inches(12.68), Inches(0.34))
    set_shape_text(title, "Jeopardy", font_size=18, bold=True, color=COLOR_ACCENT)

    left_margin = Inches(0.28)
    top_margin = Inches(0.60)
    board_width = Inches(12.77)
    board_height = Inches(6.48)
    header_height = Inches(1.02)
    row_height = (board_height - header_height) / max(1, len(values))
    col_width = board_width / max(1, len(categories))

    for col_index, category in enumerate(categories):
        left = left_margin + col_width * col_index
        header = add_rect(slide, left, top_margin, col_width, header_height, fill_color=COLOR_HEADER, line_color=COLOR_TEXT, line_width=1.4)
        set_shape_text(header, category.upper(), font_size=fit_font_size(category, max_chars_per_line=16, max_font=19, min_font=12), bold=True, color=COLOR_TEXT)

    for row_index, value in enumerate(values):
        for col_index, category in enumerate(categories):
            record = records_by_key.get((category, value))
            left = left_margin + col_width * col_index
            top = top_margin + header_height + row_height * row_index
            if record is None:
                empty_tile = add_rect(slide, left, top, col_width, row_height, fill_color=COLOR_BOARD_USED, line_color=COLOR_TEXT, line_width=1.0)
                empty_tile.fill.transparency = 0.15
                continue
            used_tile = add_rect(slide, left, top, col_width, row_height, fill_color=COLOR_BOARD_USED, line_color=COLOR_TEXT, line_width=1.0)
            set_shape_name(used_tile, f"used__{record.safe_id}")
            set_shape_text(used_tile, "", font_size=24, bold=True, color=COLOR_TEXT)
            active_tile = add_rect(slide, left, top, col_width, row_height, fill_color=COLOR_BOARD, line_color=COLOR_TEXT, line_width=1.0)
            set_shape_name(active_tile, f"tile__{record.safe_id}")
            set_shape_text(active_tile, record.value_label, font_size=25, bold=True, color=COLOR_BUTTON)


def create_clue_slide(prs: Presentation, record: ClueRecord) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    category_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.22), Inches(12.2), Inches(0.4))
    set_shape_text(category_box, record.category.upper(), font_size=22, bold=True, color=COLOR_ACCENT)
    value_box = slide.shapes.add_textbox(Inches(5.25), Inches(0.68), Inches(2.8), Inches(0.4))
    set_shape_text(value_box, record.value_label, font_size=20, bold=True, color=COLOR_BUTTON)
    panel = add_rect(slide, Inches(0.82), Inches(1.25), Inches(11.7), Inches(4.35), fill_color=COLOR_PANEL, line_color=COLOR_TEXT, line_width=1.5)
    set_shape_text(panel, record.clue, font_size=fit_font_size(record.clue, max_chars_per_line=42, max_font=26, min_font=18), bold=True, color=COLOR_TEXT)
    button = add_rect(slide, Inches(4.1), Inches(6.0), Inches(5.1), Inches(0.62), fill_color=COLOR_BUTTON, line_color=COLOR_TEXT, line_width=1.2)
    set_shape_name(button, f"reveal__{record.safe_id}")
    set_shape_text(button, "Reveal Answer", font_size=20, bold=True, color=COLOR_BG)


def create_answer_slide(prs: Presentation, record: ClueRecord) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    header = slide.shapes.add_textbox(Inches(0.55), Inches(0.22), Inches(12.2), Inches(0.4))
    set_shape_text(header, f"{record.category.upper()}  |  {record.value_label}", font_size=22, bold=True, color=COLOR_ACCENT)
    label = slide.shapes.add_textbox(Inches(4.55), Inches(1.1), Inches(4.2), Inches(0.45))
    set_shape_text(label, "ANSWER", font_size=28, bold=True, color=COLOR_BUTTON)
    panel = add_rect(slide, Inches(0.9), Inches(1.85), Inches(11.52), Inches(2.95), fill_color=COLOR_BOARD, line_color=COLOR_TEXT, line_width=1.5)
    set_shape_text(panel, record.answer, font_size=fit_font_size(record.answer, max_chars_per_line=36, max_font=28, min_font=20), bold=True, color=COLOR_TEXT)
    button = add_rect(slide, Inches(3.82), Inches(5.7), Inches(5.7), Inches(0.66), fill_color=COLOR_BUTTON, line_color=COLOR_TEXT, line_width=1.2)
    set_shape_name(button, f"return__{record.safe_id}")
    set_shape_text(button, "Return to Board", font_size=20, bold=True, color=COLOR_BG)


def create_game_over_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    title = slide.shapes.add_textbox(Inches(2.2), Inches(2.3), Inches(8.9), Inches(0.8))
    set_shape_text(title, "Game Over", font_size=34, bold=True, color=COLOR_BUTTON)
    subtitle = slide.shapes.add_textbox(Inches(1.9), Inches(3.3), Inches(9.6), Inches(0.6))
    set_shape_text(subtitle, "All clues have been played.", font_size=22, bold=False, color=COLOR_TEXT)


def build_base_presentation(records: list[ClueRecord], base_path: Path) -> dict[str, object]:
    categories = order_categories(records)
    values = order_values(records)
    records_by_key = {(record.category, record.value_label): record for record in records}

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    create_board_slide(prs, categories, values, records_by_key)

    slide_map: dict[str, SlideMeta] = {}
    for record in records:
        create_clue_slide(prs, record)
        create_answer_slide(prs, record)
        slide_map[record.safe_id] = SlideMeta(clue_index=2 * record.ordinal, answer_index=2 * record.ordinal + 1)

    create_game_over_slide(prs)
    base_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(base_path)
    return {"records": records, "slide_map": slide_map, "board_index": 1, "game_over_index": len(prs.slides), "base_path": base_path}


def generate_vba_code(meta: dict[str, object]) -> str:
    records: list[ClueRecord] = meta["records"]
    slide_map: dict[str, SlideMeta] = meta["slide_map"]
    board_index = meta["board_index"]
    game_over_index = meta["game_over_index"]

    refresh_lines = [f'    ActivePresentation.Slides({board_index}).Shapes("tile__{record.safe_id}").Visible = Not usedFlags({record.ordinal})' for record in records]
    goto_macros = []
    reveal_macros = []
    complete_macros = []
    for record in records:
        slides = slide_map[record.safe_id]
        goto_macros.append(f'Public Sub GoToClue_{record.safe_id}()\n    EnsureInitialized\n    GoToSlideIndex {slides.clue_index}\nEnd Sub')
        reveal_macros.append(f'Public Sub Reveal_{record.safe_id}()\n    EnsureInitialized\n    GoToSlideIndex {slides.answer_index}\nEnd Sub')
        complete_macros.append(
            f'Public Sub Complete_{record.safe_id}()\n'
            f'    EnsureInitialized\n'
            f'    usedFlags({record.ordinal}) = True\n'
            f'    RefreshBoard\n'
            f'    If AllCluesUsed() Then\n'
            f'        GoToSlideIndex {game_over_index}\n'
            f'    Else\n'
            f'        GoToSlideIndex {board_index}\n'
            f'    End If\n'
            f'End Sub'
        )

    return textwrap.dedent(
        f"""
        Option Explicit

        Private initialized As Boolean
        Private usedFlags(1 To {len(records)}) As Boolean

        Public Sub ResetGameState()
            Dim i As Long
            For i = 1 To {len(records)}
                usedFlags(i) = False
            Next i
            initialized = True
            RefreshBoard
        End Sub

        Private Sub EnsureInitialized()
            If Not initialized Then
                ResetGameState
            End If
        End Sub

        Private Sub RefreshBoard()
        {chr(10).join(refresh_lines)}
        End Sub

        Private Function AllCluesUsed() As Boolean
            Dim i As Long
            For i = 1 To {len(records)}
                If Not usedFlags(i) Then
                    AllCluesUsed = False
                    Exit Function
                End If
            Next i
            AllCluesUsed = True
        End Function

        Private Sub GoToSlideIndex(ByVal slideIndex As Long)
            If SlideShowWindows.Count > 0 Then
                SlideShowWindows(1).View.GotoSlide slideIndex
            End If
        End Sub

        {chr(10).join(goto_macros)}

        {chr(10).join(reveal_macros)}

        {chr(10).join(complete_macros)}
        """
    ).strip() + "\n"


def run_osascript(script: str) -> str:
    proc = subprocess.run(["osascript"], input=script, text=True, capture_output=True, check=True)
    return proc.stdout


def assign_macro_action(slide_index: int, shape_name: str, macro_name: str) -> None:
    script = f'''
tell application "Microsoft PowerPoint"
    tell active presentation
        tell slide {slide_index}
            set ast to get action setting for shape "{shape_name}" event mouse activation mouse click
            set action of ast to action type run macro
            set action setting to run of ast to "{macro_name}"
        end tell
        save
    end tell
end tell
'''
    run_osascript(script)


def inject_vba_and_save(meta: dict[str, object], output_pptm: Path, vba_code: str) -> None:
    base_pptx: Path = meta["base_path"]
    records: list[ClueRecord] = meta["records"]
    slide_map: dict[str, SlideMeta] = meta["slide_map"]

    if shutil.which("osascript") is None:
        raise RuntimeError("This script requires macOS PowerPoint automation via osascript.")

    subprocess.run(["open", "-a", "Microsoft PowerPoint", str(base_pptx)], check=True)
    run_osascript('delay 3')
    run_osascript(f'tell application "Microsoft PowerPoint" to save active presentation in POSIX file "{output_pptm}" as save as Open XML presentation macro enabled')

    vba_window_name = output_pptm.stem
    jxa_script = f"""
var se = Application('System Events');
var proc = se.processes.byName('Microsoft PowerPoint');
proc.frontmost = true;
function clickMenuBarItem(menuBarItemName, menuItemPath) {{
  var menu = proc.menuBars[0].menuBarItems.byName(menuBarItemName).menus[0];
  var item = menu;
  for (var i = 0; i < menuItemPath.length; i++) {{
    item = item.menuItems.byName(menuItemPath[i]);
    if (i < menuItemPath.length - 1) {{
      item = item.menus[0];
    }}
  }}
  item.click();
}}
clickMenuBarItem('Tools', ['Macro', 'Visual Basic Editor']);
delay(1);
var win = proc.windows.byName('Microsoft Visual Basic - {vba_window_name}');
win.attributes.byName('AXMain').value = true;
delay(0.5);
win.toolbars()[4].menuButtons.byName('Insert Module').click();
delay(0.7);
"""
    subprocess.run(["osascript", "-l", "JavaScript"], input=jxa_script, text=True, check=True)

    subprocess.run(["pbcopy"], input=vba_code, text=True, check=True)
    paste_script = """
var se = Application('System Events');
var proc = se.processes.byName('Microsoft PowerPoint');
proc.frontmost = true;
delay(0.5);
se.keystroke('a', {using: ['command down']});
delay(0.2);
se.keyCode(51);
delay(0.2);
se.keystroke('v', {using: ['command down']});
delay(1);
se.keystroke('s', {using: ['command down']});
"""
    subprocess.run(["osascript", "-l", "JavaScript"], input=paste_script, text=True, check=True)

    disable_advance = f'''
tell application "Microsoft PowerPoint"
    tell active presentation
        repeat with i from 1 to {meta["game_over_index"]}
            set advance on click of slide show transition of slide i to false
        end repeat
        save
    end tell
end tell
'''
    run_osascript(disable_advance)

    for record in records:
        slides = slide_map[record.safe_id]
        assign_macro_action(1, f"tile__{record.safe_id}", f"GoToClue_{record.safe_id}")
        assign_macro_action(slides.clue_index, f"reveal__{record.safe_id}", f"Reveal_{record.safe_id}")
        assign_macro_action(slides.answer_index, f"return__{record.safe_id}", f"Complete_{record.safe_id}")

    run_osascript('tell application "Microsoft PowerPoint" to save active presentation')
    run_osascript('tell application "Microsoft PowerPoint" to close active presentation saving yes')
    run_osascript('tell application "Microsoft PowerPoint" to quit saving yes')


def build_report(meta: dict[str, object], output_pptm: Path, report_path: Path, source_name: str) -> str:
    records: list[ClueRecord] = meta["records"]
    slide_map: dict[str, SlideMeta] = meta["slide_map"]
    with zipfile.ZipFile(output_pptm) as zf:
        has_vba = "ppt/vbaProject.bin" in zf.namelist()

    lines = [
        "Diagnosis",
        "---------",
        "A plain .pptx cannot store persistent arbitrary-order cumulative board state during slideshow runtime.",
        "This build uses VBA in a .pptm so used clues can be tracked in memory and the single board slide can be refreshed cumulatively.",
        "",
        "Architecture",
        "------------",
        f"- source: {source_name}",
        f"- total clues: {len(records)}",
        f"- slide count: {meta['game_over_index']}",
        "- deck model: one board slide, clue slides, answer slides, one Game Over slide",
        "- cumulative state: VBA usedFlags array",
        f"- vba project present: {'YES' if has_vba else 'NO'}",
        "",
        "Slide Map",
        "---------",
    ]
    for record in records:
        slides = slide_map[record.safe_id]
        lines.append(f"- {record.category} | {record.value_label} | {record.safe_id} | clue slide {slides.clue_index} | answer slide {slides.answer_index}")

    lines.extend([
        "",
        "Validation",
        "----------",
        f"- output file exists: {'YES' if output_pptm.exists() else 'NO'}",
        f"- output extension: {output_pptm.suffix}",
        f"- macro container present: {'YES' if has_vba else 'NO'}",
    ])
    report = "\n".join(lines) + "\n"
    report_path.parent.mkdir(parents=True, exist_ok=True)
    report_path.write_text(report, encoding="utf-8")
    return report


def main() -> None:
    args = parse_args()
    input_path = Path(args.input).resolve()
    output_path = Path(args.output).resolve()
    report_path = Path(args.report).resolve()
    base_path = output_path.with_name(output_path.stem + "_base.pptx")

    if not input_path.exists():
        raise FileNotFoundError(f"Input file not found: {input_path}")

    records, source_name = load_records(input_path)
    meta = build_base_presentation(records, base_path)
    vba_code = generate_vba_code(meta)
    inject_vba_and_save(meta, output_path, vba_code)
    report = build_report(meta, output_path, report_path, source_name)
    print(f"Built base deck: {base_path}")
    print(f"Built macro-enabled deck: {output_path}")
    print(f"Wrote report: {report_path}")
    print(report)


if __name__ == "__main__":
    main()
